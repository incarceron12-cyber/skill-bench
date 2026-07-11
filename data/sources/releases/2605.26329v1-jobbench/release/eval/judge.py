#!/usr/bin/env python3
"""
Generic JobBench LLM-as-Judge.

Reads model output files, evaluates them against task rubrics, and writes:
  - a scalar reward file
  - an optional detailed JSON report

The API client is OpenAI-compatible by construction, so the same judge can be
pointed at OpenAI or any compatible proxy simply by changing
JUDGE_API_BASE / JUDGE_API_KEY / JUDGE_MODEL.
"""

from __future__ import annotations

import argparse
import base64
import binascii
import fcntl
import hashlib
import json
import os
import re
import tempfile
import time
import traceback
from contextlib import contextmanager
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone
from pathlib import Path
from zipfile import BadZipFile, ZipFile


MAX_CHARS_PER_FILE = 200_000
SQLITE_EXTS = {"db", "sqlite", "sqlite3"}
SQLITE_ROWS_PER_TABLE = 500
DEFAULT_JUDGE_API_BASE = "https://api.x.ai/v1"

VISION_IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp"}
MAX_VISION_IMAGES = 8
VISUAL_RUBRIC_PATTERN = re.compile(
    r"\b(plot|figure|visualization|visualisation|visualize|visualise|"
    r"heatmap|histogram|scatter ?plot|biplot|diagram|q[- ]?q)\b",
    re.IGNORECASE,
)
VISION_MIME = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg", "gif": "image/gif", "webp": "image/webp"}


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return f"[ERROR: Failed to read text file: {path.name}: {exc}]"


def convert_file_to_text(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")

    if ext in (
        "txt", "md", "csv", "py", "json", "sh", "log", "xml", "html",
        "css", "js", "ts", "yaml", "yml", "ini", "cfg", "conf", "sql",
        "rules", "geojson",
    ):
        return _read_text(path)

    if ext in ("xlsx", "xls"):
        try:
            import pandas as pd

            xl = pd.ExcelFile(str(path))
            parts = []
            for sheet in xl.sheet_names:
                df = pd.read_excel(xl, sheet_name=sheet)
                parts.append(f"=== Sheet: {sheet} ===\n{df.to_csv(index=False)}")
            return "\n".join(parts)
        except ImportError:
            return f"[ERROR: pandas/openpyxl not available for {path.name}]"
        except Exception as exc:
            return f"[ERROR: Failed to read Excel {path.name}: {exc}]"

    if ext == "docx":
        try:
            import mammoth

            def embedded_image_placeholder(image):
                return {
                    "src": "embedded-image",
                    "alt": f"Embedded image: {image.content_type}",
                }

            with open(str(path), "rb") as f:
                result = mammoth.convert_to_markdown(
                    f,
                    convert_image=mammoth.images.img_element(embedded_image_placeholder),
                )
            return result.value
        except ImportError:
            return f"[ERROR: mammoth not available for {path.name}]"
        except Exception as exc:
            return f"[ERROR: Failed to read DOCX {path.name}: {exc}]"

    if ext == "pdf":
        try:
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                parts = []
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text(layout=True) or ""
                    parts.append(f"=== Page {i + 1} ===\n{text}")
            return "\n".join(parts)
        except ImportError:
            return f"[ERROR: pdfplumber not available for {path.name}]"
        except Exception as exc:
            return f"[ERROR: Failed to read PDF {path.name}: {exc}]"

    if ext in ("db", "sqlite", "sqlite3"):
        import sqlite3 as sqlite

        try:
            con = sqlite.connect(str(path))
            cur = con.cursor()
            schema = con.execute("SELECT sql FROM sqlite_master WHERE sql IS NOT NULL").fetchall()
            parts = ["=== Schema ==="]
            parts.extend(row[0] for row in schema if row[0])
            tables = [row[0] for row in con.execute("SELECT name FROM sqlite_master WHERE type='table'").fetchall()]
            for table in tables:
                parts.append(f"\n=== Table: {table} ===")
                try:
                    total_rows = con.execute(f'SELECT COUNT(*) FROM "{table}"').fetchone()[0]
                    rows = cur.execute(f'SELECT * FROM "{table}" LIMIT {SQLITE_ROWS_PER_TABLE}').fetchall()
                    cols = [d[0] for d in cur.description]
                    parts.append(f"-- total_rows: {total_rows}; shown: {len(rows)} (LIMIT {SQLITE_ROWS_PER_TABLE})")
                    parts.append(",".join(cols))
                    for row in rows:
                        parts.append(",".join("" if value is None else str(value) for value in row))
                except Exception as exc:
                    parts.append(f"[ERROR reading table {table}: {exc}]")
            con.close()
            return "\n".join(parts)
        except Exception as exc:
            return f"[ERROR: Failed to read SQLite {path.name}: {exc}]"

    if ext == "pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            parts = []
            for idx, slide in enumerate(prs.slides):
                parts.append(f"=== Slide {idx + 1} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts)
        except ImportError:
            return f"[ERROR: python-pptx not available for {path.name}]"
        except Exception as exc:
            return f"[ERROR: Failed to read PowerPoint {path.name}: {exc}]"

    if ext == "ipynb":
        try:
            nb = json.loads(path.read_text(encoding="utf-8"))
            parts = []
            for cell in nb.get("cells", []):
                parts.append(f"=== {cell['cell_type']} ===")
                parts.append("".join(cell.get("source", [])))
                for output in cell.get("outputs", []):
                    if "text" in output:
                        parts.append("".join(output["text"]))
            return "\n".join(parts)
        except Exception as exc:
            return f"[ERROR: Failed to read notebook {path.name}: {exc}]"

    if ext in ("png", "jpg", "jpeg", "gif", "svg", "bmp"):
        return f"[Image file: {path.name} — cannot extract text content]"

    return f"[Binary or unsupported file type: {ext} — {path.name}]"


def extract_all_file_contents(output_dir: Path) -> str:
    parts = []
    for file_path in sorted(output_dir.rglob("*")):
        if not file_path.is_file():
            continue
        content = convert_file_to_text(file_path)
        ext = file_path.suffix.lower().lstrip(".")
        if ext not in SQLITE_EXTS and len(content) > MAX_CHARS_PER_FILE:
            content = content[:MAX_CHARS_PER_FILE] + f"\n... [Content truncated at {MAX_CHARS_PER_FILE} characters]"
        parts.append(f"=== FILE: {file_path.name} ===\n{content}\n")
    return "\n".join(parts)


def rubric_needs_vision(rubric: dict) -> bool:
    text = rubric.get("rubric", "") or ""
    criterion = rubric.get("criterion", [])
    if isinstance(criterion, list):
        text = text + " " + " ".join(criterion)
    elif isinstance(criterion, str):
        text = text + " " + criterion
    return bool(VISUAL_RUBRIC_PATTERN.search(text))


def collect_image_paths(output_dir: Path, cap: int | None = MAX_VISION_IMAGES) -> list[Path]:
    if not output_dir.exists():
        return []
    images = [
        p
        for p in sorted(output_dir.rglob("*"))
        if p.is_file() and p.suffix.lower().lstrip(".") in VISION_IMAGE_EXTS
    ]
    return images if cap is None else images[:cap]


def image_to_data_url(path: Path) -> str | None:
    ext = path.suffix.lower().lstrip(".")
    mime = VISION_MIME.get(ext)
    if mime is None:
        return None
    try:
        encoded = base64.b64encode(path.read_bytes()).decode("ascii")
    except Exception:
        return None
    return f"data:{mime};base64,{encoded}"


def collect_image_attachments(
    output_dir: Path,
    cap: int = MAX_VISION_IMAGES,
) -> list[tuple[str, str]]:
    """Collect standalone and embedded images as deduplicated data URLs."""
    if not output_dir.exists() or cap <= 0:
        return []

    attachments: list[tuple[str, str]] = []
    seen_hashes: set[str] = set()

    def add_attachment(name: str, mime: str, image_bytes: bytes) -> bool:
        digest = hashlib.sha256(image_bytes).hexdigest()
        if digest in seen_hashes:
            return False
        seen_hashes.add(digest)
        encoded = base64.b64encode(image_bytes).decode("ascii")
        attachments.append((name, f"data:{mime};base64,{encoded}"))
        return len(attachments) >= cap

    for path in collect_image_paths(output_dir, cap=None):
        ext = path.suffix.lower().lstrip(".")
        mime = VISION_MIME.get(ext)
        if mime is None:
            continue
        try:
            image_bytes = path.read_bytes()
        except OSError:
            continue
        if add_attachment(path.relative_to(output_dir).as_posix(), mime, image_bytes):
            return attachments

    for docx_path in sorted(output_dir.rglob("*.docx")):
        try:
            with ZipFile(docx_path) as archive:
                media_names = [
                    name
                    for name in sorted(archive.namelist())
                    if name.startswith("word/media/")
                    and Path(name).suffix.lower().lstrip(".") in VISION_IMAGE_EXTS
                ]
                for media_name in media_names:
                    ext = Path(media_name).suffix.lower().lstrip(".")
                    mime = VISION_MIME.get(ext)
                    if mime is None:
                        continue
                    try:
                        image_bytes = archive.read(media_name)
                    except (KeyError, OSError):
                        continue
                    display_name = f"{docx_path.relative_to(output_dir).as_posix()}:{media_name}"
                    if add_attachment(display_name, mime, image_bytes):
                        return attachments
        except (BadZipFile, OSError):
            continue

    for notebook_path in sorted(output_dir.rglob("*.ipynb")):
        try:
            notebook = json.loads(notebook_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            continue
        for cell_index, cell in enumerate(notebook.get("cells", [])):
            for output_index, output in enumerate(cell.get("outputs", [])):
                data = output.get("data", {})
                if not isinstance(data, dict):
                    continue
                for mime in ("image/png", "image/jpeg", "image/gif", "image/webp"):
                    encoded = data.get(mime)
                    if isinstance(encoded, list):
                        encoded = "".join(str(part) for part in encoded)
                    if not isinstance(encoded, str):
                        continue
                    try:
                        image_bytes = base64.b64decode("".join(encoded.split()), validate=True)
                    except (ValueError, binascii.Error):
                        continue
                    display_name = (
                        f"{notebook_path.relative_to(output_dir).as_posix()}:"
                        f"cell-{cell_index}-output-{output_index}:{mime}"
                    )
                    if add_attachment(display_name, mime, image_bytes):
                        return attachments

    return attachments


def normalize_criteria(rubric: dict) -> list[str]:
    criterion_raw = rubric.get("criterion", [])
    if isinstance(criterion_raw, str):
        return [criterion_raw]
    return list(criterion_raw)


def build_failed_rubric_result(
    rubric_index: int,
    rubric: dict,
    overall_reasoning: str,
    criteria_reasoning: str | None = None,
) -> dict:
    criteria = normalize_criteria(rubric)
    per_criterion_reasoning = criteria_reasoning or overall_reasoning
    return {
        "index": rubric_index,
        "rubric": rubric.get("rubric", ""),
        "weight": rubric.get("weight", 0),
        "result": {
            "passed": False,
            "score": 0,
            "criteria_count": len(criteria),
            "criteria_passed": 0,
            "criteria_results": [
                {
                    "index": idx,
                    "criterion": criterion,
                    "passed": False,
                    "reasoning": per_criterion_reasoning,
                    "evidence": "",
                }
                for idx, criterion in enumerate(criteria)
            ],
            "overall_reasoning": overall_reasoning,
        },
    }


def build_scorecard(results: list[dict]) -> dict[str, float | int]:
    total_score = sum(result["result"]["score"] for result in results)
    max_score = sum(result["weight"] for result in results)
    passed_count = sum(1 for result in results if result["result"]["passed"])
    total_count = len(results)
    normalized = round(total_score / max_score, 4) if max_score > 0 else 0.0
    pass_rate = round(passed_count / total_count, 4) if total_count > 0 else 0.0

    scorecard: dict[str, float | int] = {
        "total_score": total_score,
        "max_score": max_score,
        "normalized_score": normalized,
        "pass_rate": pass_rate,
        "passed_count": passed_count,
        "total_count": total_count,
    }
    for result in results:
        idx = result["index"]
        scorecard[f"rubric_{idx}_passed"] = 1 if result["result"]["passed"] else 0
        scorecard[f"rubric_{idx}_score"] = result["result"]["score"]
    return scorecard


def build_reward(scorecard: dict[str, float | int]) -> dict[str, float]:
    return {"reward": float(scorecard.get("normalized_score", 0.0))}


def build_details_report(
    evaluated_model: str,
    judge_model: str,
    results: list[dict],
    total_count: int | None = None,
) -> dict:
    total_score = sum(result["result"]["score"] for result in results)
    max_score = sum(result["weight"] for result in results)
    passed_count = sum(1 for result in results if result["result"]["passed"])
    effective_total_count = total_count if total_count is not None else len(results)
    pass_rate_value = int((passed_count / effective_total_count) * 100) if effective_total_count > 0 else 0

    return {
        "evaluated_model": evaluated_model,
        "judge_model": judge_model,
        "timestamp": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "total_score": total_score,
        "max_score": max_score,
        "pass_rate": f"{pass_rate_value}%",
        "passed_count": passed_count,
        "total_count": effective_total_count,
        "rubrics": results,
    }


def write_json(path: Path, payload: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(
        "w",
        encoding="utf-8",
        delete=False,
        dir=path.parent,
        prefix=f".{path.name}.",
        suffix=".tmp",
    ) as handle:
        json.dump(payload, handle, ensure_ascii=False, indent=2)
        handle.write("\n")
        tmp_path = Path(handle.name)
    os.replace(tmp_path, path)


def write_outputs(
    result_file: Path | None,
    reward: dict[str, float],
    details_file: Path | None,
    details: dict | None,
) -> None:
    if result_file is not None:
        write_json(result_file, reward)
    if details_file is not None:
        write_json(details_file, details or {})


@contextmanager
def file_lock(lock_file: Path):
    lock_file.parent.mkdir(parents=True, exist_ok=True)
    with lock_file.open("a+", encoding="utf-8") as handle:
        fcntl.flock(handle.fileno(), fcntl.LOCK_EX)
        try:
            yield
        finally:
            fcntl.flock(handle.fileno(), fcntl.LOCK_UN)


def load_existing_details(details_file: Path | None) -> dict | None:
    if details_file is None or not details_file.exists():
        return None
    try:
        payload = json.loads(details_file.read_text(encoding="utf-8"))
    except Exception:
        return None
    if isinstance(payload, dict):
        return payload
    return None


def load_existing_rubric_results(details_file: Path | None) -> dict[int, dict]:
    payload = load_existing_details(details_file)
    if not payload:
        return {}

    results: dict[int, dict] = {}
    for rubric_result in payload.get("rubrics", []):
        if not isinstance(rubric_result, dict):
            continue
        index = rubric_result.get("index")
        if isinstance(index, int):
            results[index] = rubric_result
    return results


def append_rubric_result(
    details_file: Path,
    lock_file: Path,
    evaluated_model: str,
    judge_model: str,
    total_count: int,
    rubric_result: dict,
) -> None:
    with file_lock(lock_file):
        current_payload = load_existing_details(details_file) or {}
        current_rubrics = [
            item
            for item in current_payload.get("rubrics", [])
            if isinstance(item, dict) and item.get("index") != rubric_result.get("index")
        ]
        current_rubrics.append(rubric_result)
        current_rubrics.sort(key=lambda item: item.get("index", -1))
        updated_payload = build_details_report(
            evaluated_model,
            judge_model,
            current_rubrics,
            total_count=total_count,
        )
        write_json(details_file, updated_payload)


def first_non_empty(*values: str | None) -> str:
    for value in values:
        if value:
            return value
    return ""


def resolve_api_config(
    judge_model: str,
    api_base_arg: str | None,
    api_key_arg: str | None,
) -> tuple[str, str]:
    if not judge_model:
        raise ValueError("No judge model provided. Set --judge-model or JUDGE_MODEL.")

    api_base = first_non_empty(api_base_arg, os.environ.get("JUDGE_API_BASE"), DEFAULT_JUDGE_API_BASE)
    api_key = first_non_empty(api_key_arg, os.environ.get("JUDGE_API_KEY"))
    return api_base, api_key


def get_openai_client(api_base: str, api_key: str):
    from openai import OpenAI

    if not api_key:
        raise ValueError("No judge API key provided. Set JUDGE_API_KEY.")

    return OpenAI(base_url=api_base, api_key=api_key)


def parse_judge_json(content: str) -> tuple[dict, str]:
    try:
        return json.loads(content), "direct_json"
    except json.JSONDecodeError:
        pass

    fence = re.search(r"```(?:json)?\s*\n(.*?)\n\s*```", content, re.DOTALL)
    if fence:
        try:
            return json.loads(fence.group(1).strip()), "markdown_fence"
        except json.JSONDecodeError:
            pass

    first = content.find("{")
    last = content.rfind("}")
    if first != -1 and last > first:
        try:
            return json.loads(content[first:last + 1]), "first_last_brace"
        except json.JSONDecodeError:
            pass

    for candidate in reversed(re.findall(r"\{.*?\"criteria_results\"\s*:\s*\[.*?\].*?\}", content, re.DOTALL)):
        try:
            return json.loads(candidate), "regex_extract"
        except json.JSONDecodeError:
            continue

    raise ValueError(f"Could not extract JSON from response: {content[:500]}")


def judge_rubric(
    rubric_index: int,
    rubric: dict,
    file_contents: str,
    judge_model: str,
    api_base: str,
    api_key: str,
    timeout_sec: int = 300,
    max_retries: int = 1,
    image_attachments: list[tuple[str, str]] | None = None,
) -> tuple[dict, dict]:
    rubric_text = rubric.get("rubric", "")
    weight = rubric.get("weight", 0)
    criteria = normalize_criteria(rubric)
    criterion_count = len(criteria)
    criteria_list_text = "\n".join(f"Criterion {idx}: {criterion}" for idx, criterion in enumerate(criteria))

    attached_images: list[tuple[str, str]] = []
    if image_attachments and rubric_needs_vision(rubric):
        attached_images.extend(image_attachments)
    vision_used = bool(attached_images)

    prompt = f"""You are an evaluation judge. Your task is to evaluate ALL criteria for a single rubric.

## Rubric Description
{rubric_text}

## Criteria to Evaluate (Judge ALL of them)
{criteria_list_text}

## Output Files Content
The following are the contents of all output files to evaluate:

{file_contents}

## Evaluation Rules
- Evaluate EACH criterion listed above independently
- For each criterion: determine if it PASSES or FAILS
- Semantic matching is acceptable (you don't need exact wording match)
- Binary judgment for each criterion: PASS or FAIL only
- The rubric passes ONLY if ALL criteria pass

## Output Format
Return your judgment as a JSON object with EXACTLY this structure (no markdown, no extra text):
{{
  "criteria_results": [
    {{"index": 0, "passed": true/false, "reasoning": "...", "evidence": "..."}},
    {{"index": 1, "passed": true/false, "reasoning": "...", "evidence": "..."}}
  ],
  "rubric_passed": true/false,
  "overall_reasoning": "Summary of why the rubric passed or failed"
}}

IMPORTANT:
- criteria_results array must have exactly {criterion_count} items (one for each criterion)
- rubric_passed should be true ONLY if ALL criteria passed
- Include specific evidence from the output files{' and the attached images' if vision_used else ''}
"""

    if vision_used:
        user_content: list[dict] = [{"type": "text", "text": prompt}]
        user_content.append({
            "type": "text",
            "text": f"\n## Attached Images ({len(attached_images)} file{'s' if len(attached_images) != 1 else ''})",
        })
        for i, (fname, url) in enumerate(attached_images, start=1):
            user_content.append({"type": "text", "text": f"Image {i}: {fname}"})
            user_content.append({"type": "image_url", "image_url": {"url": url}})
    else:
        user_content = prompt

    last_error = None
    raw_response = ""
    parse_status = "failed"
    for attempt in range(max_retries):
        try:
            client = get_openai_client(api_base, api_key)
            response = client.chat.completions.create(
                model=judge_model,
                messages=[
                    {
                        "role": "system",
                        "content": "You are an evaluation judge. You must return valid JSON only, with no markdown formatting or extra text.",
                    },
                    {"role": "user", "content": user_content},
                ],
                max_completion_tokens=8192,
                temperature=0.0,
                timeout=timeout_sec,
            )
            raw_response = response.choices[0].message.content.strip()
            parsed, parse_status = parse_judge_json(raw_response)

            model_criteria = parsed.get("criteria_results", [])
            rubric_passed = bool(parsed.get("rubric_passed", False))
            overall_reasoning = parsed.get("overall_reasoning", "")

            enriched = []
            for idx, criterion in enumerate(criteria):
                item = model_criteria[idx] if idx < len(model_criteria) else {}
                enriched.append(
                    {
                        "index": idx,
                        "criterion": criterion,
                        "passed": bool(item.get("passed", False)),
                        "reasoning": item.get("reasoning", ""),
                        "evidence": item.get("evidence", ""),
                    }
                )

            score = weight if rubric_passed else 0
            criteria_passed = sum(1 for item in enriched if item["passed"])
            result = {
                "index": rubric_index,
                "rubric": rubric_text,
                "weight": weight,
                "result": {
                    "passed": rubric_passed,
                    "score": score,
                    "criteria_count": criterion_count,
                    "criteria_passed": criteria_passed,
                    "criteria_results": enriched,
                    "overall_reasoning": overall_reasoning,
                },
            }
            debug = {
                "api_base": api_base,
                "parse_status": parse_status,
                "api_exit_code": 0,
                "criterion_count": criterion_count,
                "criteria_list_text": criteria_list_text,
                "rubric_text": rubric_text,
                "raw_response": raw_response,
                "vision_used": vision_used,
                "attached_images": [name for name, _ in attached_images],
            }
            return result, debug
        except Exception as exc:
            last_error = exc
            if attempt < max_retries - 1:
                time.sleep(2)
            continue

    default_criteria = [
        {
            "index": idx,
            "criterion": criterion,
            "passed": False,
            "reasoning": f"Failed to get judge response: {last_error}",
            "evidence": "",
        }
        for idx, criterion in enumerate(criteria)
    ]
    result = {
        "index": rubric_index,
        "rubric": rubric_text,
        "weight": weight,
        "result": {
            "passed": False,
            "score": 0,
            "criteria_count": criterion_count,
            "criteria_passed": 0,
            "criteria_results": default_criteria,
            "overall_reasoning": f"Failed after {max_retries} retries: {last_error}",
        },
    }
    debug = {
        "api_base": api_base,
        "parse_status": parse_status,
        "api_exit_code": 1 if raw_response else 2,
        "criterion_count": criterion_count,
        "criteria_list_text": criteria_list_text,
        "rubric_text": rubric_text,
        "raw_response": raw_response,
        "error": str(last_error) if last_error is not None else "",
        "vision_used": vision_used,
        "attached_images": [name for name, _ in attached_images],
    }
    return result, debug


def write_detail_log(
    detail_log_dir: Path | None,
    detail_log_prefix: str,
    rubric_index: int,
    judge_model: str,
    debug: dict,
    final_result: dict,
) -> None:
    if detail_log_dir is None or not detail_log_prefix:
        return

    detail_log_dir.mkdir(parents=True, exist_ok=True)
    detail_log_file = detail_log_dir / f"{detail_log_prefix}_rubric_{rubric_index}.log"
    sections = [
        "========================================",
        "Rubric Judge Detail Log",
        "========================================",
        f"Timestamp: {datetime.now(timezone.utc).isoformat(timespec='seconds')}",
        f"Unique Key: {detail_log_prefix}",
        f"Rubric Index: {rubric_index}",
        f"Judge Model: {judge_model}",
        f"API Base: {debug.get('api_base', '')}",
        f"Parse Status: {debug.get('parse_status', 'failed')}",
        f"API Exit Code: {debug.get('api_exit_code', '')}",
        f"Criteria Count: {debug.get('criterion_count', '')}",
        f"Vision Used: {debug.get('vision_used', False)}",
        f"Attached Images: {', '.join(debug.get('attached_images', [])) or '(none)'}",
    ]
    if debug.get("error"):
        sections.append(f"Error: {debug['error']}")

    sections.extend(
        [
            "",
            "========================================",
            "RUBRIC TEXT",
            "========================================",
            debug.get("rubric_text", ""),
            "",
            "========================================",
            "CRITERIA",
            "========================================",
            debug.get("criteria_list_text", ""),
            "",
            "========================================",
            "RAW API RESPONSE",
            "========================================",
            debug.get("raw_response", ""),
            "",
            "========================================",
            "FINAL RESULT",
            "========================================",
            json.dumps(final_result, ensure_ascii=False, indent=2),
            "",
            "======================================== END ========================================",
            "",
        ]
    )
    detail_log_file.write_text("\n".join(sections), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="JobBench LLM judge")
    parser.add_argument("--output-dir", required=True, help="Directory with model output files")
    parser.add_argument("--rubrics-file", required=True, help="Path to RUBRICS.json")
    parser.add_argument("--result-file", default=None, help="Optional path for reward JSON")
    parser.add_argument("--details-file", default=None, help="Where to write detailed results JSON")
    parser.add_argument("--judge-model", default=os.environ.get("JUDGE_MODEL", ""))
    parser.add_argument("--api-base", default=None, help="OpenAI-compatible API base URL")
    parser.add_argument("--api-key", default=None, help="OpenAI-compatible API key")
    parser.add_argument("--max-workers", type=int, default=10)
    parser.add_argument("--max-retries", type=int, default=1)
    parser.add_argument("--timeout-per-rubric", type=int, default=300)
    parser.add_argument("--evaluated-model", default="", help="Name of the model output being judged")
    parser.add_argument("--lock-file", default=None, help="Optional lock file path for incremental details writes")
    parser.add_argument("--detail-log-dir", default=None, help="Optional directory for per-rubric detail logs")
    parser.add_argument("--detail-log-prefix", default="", help="Prefix for per-rubric detail logs")
    args = parser.parse_args()

    output_dir = Path(args.output_dir)
    rubrics_file = Path(args.rubrics_file)
    result_file = Path(args.result_file) if args.result_file else None
    details_file = Path(args.details_file) if args.details_file else None
    lock_file = Path(args.lock_file) if args.lock_file else None
    detail_log_dir = Path(args.detail_log_dir) if args.detail_log_dir else None
    if lock_file is None and details_file is not None:
        lock_file = details_file.with_name(f".{details_file.stem}.lock")

    try:
        if not rubrics_file.exists():
            write_outputs(
                result_file,
                build_reward(build_scorecard([])),
                details_file,
                {"error": "rubrics not found", "rubrics_file": str(rubrics_file)},
            )
            return

        rubrics_data = json.loads(rubrics_file.read_text(encoding="utf-8"))
        rubrics = rubrics_data.get("rubrics") or rubrics_data.get("evaluation_rubrics") or []
        if not rubrics:
            write_outputs(
                result_file,
                build_reward(build_scorecard([])),
                details_file,
                {"error": "no rubrics", "rubrics_file": str(rubrics_file)},
            )
            return

        total_rubric_count = len(rubrics)
        evaluated_model = args.evaluated_model or output_dir.name
        has_output_files = output_dir.exists() and any(path.is_file() for path in output_dir.rglob("*"))
        if not has_output_files:
            results = [
                build_failed_rubric_result(idx, rubric, "No output files found in the model output directory.")
                for idx, rubric in enumerate(rubrics)
            ]
        else:
            file_contents = extract_all_file_contents(output_dir)
            if not file_contents.strip():
                results = [
                    build_failed_rubric_result(idx, rubric, "Output files were unreadable or empty after conversion.")
                    for idx, rubric in enumerate(rubrics)
                ]
            else:
                existing_results = load_existing_rubric_results(details_file)
                api_base, api_key = resolve_api_config(args.judge_model, args.api_base, args.api_key)
                image_attachments = collect_image_attachments(output_dir)

                results: list[dict | None] = [existing_results.get(idx) for idx in range(len(rubrics))]
                with ThreadPoolExecutor(max_workers=args.max_workers) as executor:
                    futures = {
                        executor.submit(
                            judge_rubric,
                            idx,
                            rubric,
                            file_contents,
                            args.judge_model,
                            api_base,
                            api_key,
                            args.timeout_per_rubric,
                            args.max_retries,
                            image_attachments,
                        ): idx
                        for idx, rubric in enumerate(rubrics)
                        if results[idx] is None
                    }
                    for future in as_completed(futures):
                        idx = futures[future]
                        try:
                            results[idx], debug = future.result()
                        except Exception as exc:
                            results[idx] = build_failed_rubric_result(
                                idx,
                                rubrics[idx],
                                f"Judge raised an exception: {exc}",
                            )
                            debug = {
                                "api_base": api_base,
                                "parse_status": "failed",
                                "api_exit_code": 2,
                                "criterion_count": len(normalize_criteria(rubrics[idx])),
                                "criteria_list_text": "\n".join(
                                    f"Criterion {criterion_idx}: {criterion}"
                                    for criterion_idx, criterion in enumerate(normalize_criteria(rubrics[idx]))
                                ),
                                "rubric_text": rubrics[idx].get("rubric", ""),
                                "raw_response": "",
                                "error": str(exc),
                                "vision_used": False,
                                "attached_images": [],
                            }
                        if details_file is not None and lock_file is not None and results[idx] is not None:
                            append_rubric_result(
                                details_file,
                                lock_file,
                                evaluated_model,
                                args.judge_model,
                                total_rubric_count,
                                results[idx],
                            )
                        write_detail_log(
                            detail_log_dir,
                            args.detail_log_prefix,
                            idx,
                            args.judge_model,
                            debug,
                            results[idx],
                        )

                results = [result for result in results if result is not None]

        scorecard = build_scorecard(results)
        reward = build_reward(scorecard)
        details = build_details_report(
            evaluated_model,
            args.judge_model,
            results,
            total_count=total_rubric_count,
        )
        write_outputs(result_file, reward, details_file, details)
    except Exception as exc:
        fallback_reward = build_reward(build_scorecard([]))
        write_outputs(
            result_file,
            fallback_reward,
            details_file,
            {
                "error": str(exc),
                "traceback": traceback.format_exc(),
                "evaluated_model": args.evaluated_model or output_dir.name,
                "judge_model": args.judge_model,
            },
        )


if __name__ == "__main__":
    main()
